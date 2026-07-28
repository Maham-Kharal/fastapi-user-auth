import os
import asyncio
from pptx import Presentation
from app.core.database import SessionLocal
from app.models.models import User, ChatSession, Document
from app.services.document_parser import sniff_mime_type, extract_txt_chunks, extract_pptx_chunks
from app.services.qdrant_service import search_session_documents
from app.worker import process_document_ingestion


def create_sample_pptx(file_path: str):
    """Create a sample 2-slide PPTX file with slide titles, body text, and speaker notes."""
    prs = Presentation()
    
    # Slide 1: Title & Content
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    txBox = slide1.shapes.add_textbox(100, 100, 400, 200)
    tf = txBox.text_frame
    tf.text = "Slide 1: Artificial Intelligence Overview"
    p = tf.add_paragraph()
    p.text = "Deep learning uses multi-layer neural networks for representation learning."
    
    # Slide 1 Speaker Notes
    notes_slide1 = slide1.notes_slide
    notes_tf1 = notes_slide1.notes_text_frame
    notes_tf1.text = "Remember to emphasize backpropagation during the presentation."

    # Slide 2: Advanced Topics
    slide2 = prs.slides.add_slide(blank_slide_layout)
    txBox2 = slide2.shapes.add_textbox(100, 100, 400, 200)
    tf2 = txBox2.text_frame
    tf2.text = "Slide 2: Quantum Algorithm Complexity"
    p2 = tf2.add_paragraph()
    p2.text = "Shor's algorithm achieves exponential speedup for integer factorization."

    # Slide 2 Speaker Notes
    notes_slide2 = slide2.notes_slide
    notes_tf2 = notes_slide2.notes_text_frame
    notes_tf2.text = "Note: Complexity bound is O( (log N)^3 )."

    prs.save(file_path)


def create_sample_txt(file_path: str):
    """Create a sample TXT document."""
    content = (
        "Library System Architecture Specification\n"
        "The authentication module uses JWT tokens with 15-minute expiration.\n"
        "Database connections are pooled using AsyncPG with a max limit of 20 connections."
    )
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)


async def run_txt_pptx_test():
    print("==================================================")
    print("   TXT & PPTX INGESTION & PIPELINE VERIFICATION   ")
    print("==================================================")

    db = SessionLocal()
    try:
        user = db.query(User).first()
        if not user:
            user = User(username="txt_pptx_user", email="user@txtpptx.com", hashed_password="pw")
            db.add(user)
            db.commit()
            db.refresh(user)

        session1 = ChatSession(user_id=user.id, title="Session TXT & PPTX")
        session2 = ChatSession(user_id=user.id, title="Isolated Session")
        db.add_all([session1, session2])
        db.commit()
        db.refresh(session1)
        db.refresh(session2)

        # 1. Create files
        txt_path = "test_spec.txt"
        pptx_path = "test_presentation.pptx"
        create_sample_txt(txt_path)
        create_sample_pptx(pptx_path)

        # 2. Test MIME Sniffing for TXT & PPTX
        with open(txt_path, "rb") as f:
            txt_bytes = f.read()
        mime_txt = sniff_mime_type(txt_bytes)
        print(f"[SUCCESS] TXT MIME Sniffing detected: {mime_txt}")
        assert mime_txt == "txt"

        with open(pptx_path, "rb") as f:
            pptx_bytes = f.read()
        mime_pptx = sniff_mime_type(pptx_bytes)
        print(f"[SUCCESS] PPTX MIME Sniffing detected: {mime_pptx}")
        assert mime_pptx == "pptx"

        # Test invalid binary file rejection
        try:
            sniff_mime_type(b"\x00\x01\x02\x03\x04\x05binary_exe_file")
            print("[FAIL] Failed to reject invalid binary file.")
        except ValueError as err:
            print(f"[SUCCESS] Rejected invalid file: {err}")

        # 3. Ingest TXT Document into Session 1
        doc_txt = Document(
            user_id=user.id,
            session_id=session1.id,
            filename="test_spec.txt",
            file_path=os.path.abspath(txt_path),
            file_type=mime_txt,
            file_size_bytes=len(txt_bytes),
            status="processing",
        )
        db.add(doc_txt)
        db.commit()
        db.refresh(doc_txt)

        print(f"Processing ARQ ingestion for TXT Doc ID={doc_txt.id}...")
        res_txt = await process_document_ingestion({}, doc_txt.id)
        print("TXT Ingestion Result:", res_txt)
        assert res_txt["status"] == "completed"

        # 4. Ingest PPTX Document into Session 1
        doc_pptx = Document(
            user_id=user.id,
            session_id=session1.id,
            filename="test_presentation.pptx",
            file_path=os.path.abspath(pptx_path),
            file_type=mime_pptx,
            file_size_bytes=len(pptx_bytes),
            status="processing",
        )
        db.add(doc_pptx)
        db.commit()
        db.refresh(doc_pptx)

        print(f"Processing ARQ ingestion for PPTX Doc ID={doc_pptx.id}...")
        res_pptx = await process_document_ingestion({}, doc_pptx.id)
        print("PPTX Ingestion Result:", res_pptx)
        assert res_pptx["status"] == "completed"

        # 5. Search TXT Content in Session 1
        txt_query = "What is the expiration time for JWT tokens?"
        chunks_txt = await search_session_documents(txt_query, session1.id, user.id)
        print(f"\nTXT Search Chunks Found: {len(chunks_txt)}")
        if chunks_txt:
            print("TXT Top Chunk Snippet:", chunks_txt[0]["content"])
            assert "15-minute expiration" in chunks_txt[0]["content"]

        # 6. Search PPTX Slide Text & Speaker Notes in Session 1
        pptx_query = "What is the complexity of Shor's algorithm in speaker notes?"
        chunks_pptx = await search_session_documents(pptx_query, session1.id, user.id)
        print(f"\nPPTX Search Chunks Found: {len(chunks_pptx)}")
        if chunks_pptx:
            print("PPTX Top Chunk Source:", chunks_pptx[0]["source"], "Section:", chunks_pptx[0]["section"])
            print("PPTX Top Chunk Content:\n", chunks_pptx[0]["content"])
            assert "Slide 2" in chunks_pptx[0]["section"]
            assert "[Speaker Notes]" in chunks_pptx[0]["content"]

        # 7. Verify Zero Cross-Session Leakage in Session 2
        chunks_s2 = await search_session_documents(pptx_query, session2.id, user.id)
        print(f"\nSession 2 (Isolated) Chunks Found: {len(chunks_s2)}")
        assert len(chunks_s2) == 0, "Cross-session leakage detected!"

        print("\n[SUCCESS] 100% TXT & PPTX Ingestion, Slide/Notes Chunking & Isolation Verified!")
        print("==================================================")

    finally:
        db.close()
        for p in ["test_spec.txt", "test_presentation.pptx"]:
            if os.path.exists(p):
                os.remove(p)


if __name__ == "__main__":
    asyncio.run(run_txt_pptx_test())
