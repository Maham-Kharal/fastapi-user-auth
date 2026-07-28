import os
import io
import zipfile
import logging
from typing import List, Dict, Any
from pypdf import PdfReader
import docx
from pptx import Presentation
import filetype

logger = logging.getLogger(__name__)


def sniff_mime_type(content: bytes) -> str:
    """
    Inspect raw binary magic bytes to determine file type.
    Does NOT rely on client-supplied headers or filename extensions.
    Returns 'pdf', 'docx', 'pptx', or 'txt'. Raises ValueError if invalid.
    """
    if len(content) < 4:
        raise ValueError("File is too small or corrupted.")

    # 1. PDF check: Starts with %PDF- (0x25 0x50 0x44 0x46 0x2D)
    if content.startswith(b"%PDF-"):
        return "pdf"

    # 2. OpenXML check (DOCX / PPTX): Zip container starting with PK\x03\x04
    if content.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                namelist = zf.namelist()
                if any("word/document.xml" in name for name in namelist):
                    return "docx"
                if any("ppt/presentation.xml" in name for name in namelist):
                    return "pptx"
        except zipfile.BadZipFile:
            pass

    # 3. TXT check: Valid UTF-8 printable text without binary null bytes
    if b"\x00" not in content[:2048]:
        try:
            content.decode("utf-8")
            return "txt"
        except UnicodeDecodeError:
            try:
                content.decode("latin-1")
                return "txt"
            except Exception:
                pass

    # 4. Fallback to filetype library detection
    kind = filetype.guess(content)
    if kind:
        ext = kind.extension.lower()
        if ext == "pdf":
            return "pdf"
        elif ext == "docx":
            return "docx"
        elif ext == "pptx":
            return "pptx"

    raise ValueError("Invalid file format. Supported documents are PDF, DOCX, PPTX, and TXT.")


def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
    """Split text into overlapping character chunks."""
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += (chunk_size - overlap)
    return chunks


def extract_pdf_chunks(file_path: str, filename: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict[str, Any]]:
    """Extract page-by-page text from PDF file and return chunks with page metadata."""
    reader = PdfReader(file_path)
    chunks = []

    for page_num, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        text_chunks = _chunk_text(page_text, chunk_size, overlap)
        for chunk in text_chunks:
            chunks.append({
                "content": chunk,
                "page_number": page_num,
                "section": f"Page {page_num}",
                "source": filename,
            })

    logger.info("Extracted %d chunks from PDF '%s' across %d pages.", len(chunks), filename, len(reader.pages))
    return chunks


def extract_docx_chunks(file_path: str, filename: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict[str, Any]]:
    """Extract section-by-section text from DOCX file and return chunks with section metadata."""
    doc = docx.Document(file_path)
    chunks = []
    current_section = "General"
    section_text_buffer = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            if section_text_buffer:
                full_section_text = "\n".join(section_text_buffer)
                for chunk in _chunk_text(full_section_text, chunk_size, overlap):
                    chunks.append({
                        "content": chunk,
                        "page_number": 1,
                        "section": current_section,
                        "source": filename,
                    })
                section_text_buffer = []
            current_section = text
        else:
            section_text_buffer.append(text)

    if section_text_buffer:
        full_section_text = "\n".join(section_text_buffer)
        for chunk in _chunk_text(full_section_text, chunk_size, overlap):
            chunks.append({
                "content": chunk,
                "page_number": 1,
                "section": current_section,
                "source": filename,
            })

    logger.info("Extracted %d chunks from DOCX '%s'.", len(chunks), filename)
    return chunks


def extract_txt_chunks(file_path: str, filename: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict[str, Any]]:
    """Extract text from TXT file and return overlapping chunks."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text_content = f.read()

    text_chunks = _chunk_text(text_content, chunk_size, overlap)
    chunks = [
        {
            "content": chunk,
            "page_number": 1,
            "section": "General",
            "source": filename,
        }
        for chunk in text_chunks
    ]

    logger.info("Extracted %d chunks from TXT '%s'.", len(chunks), filename)
    return chunks


def extract_pptx_chunks(file_path: str, filename: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict[str, Any]]:
    """
    Extract slide-by-slide text and speaker notes from PPTX file.
    Speaker notes are attached to the slide context under '[Speaker Notes]:'.
    Preserves slide numbers for precise citations like [filename.pptx, Slide 2].
    """
    prs = Presentation(file_path)
    chunks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        
        # 1. Extract slide text elements
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)

        # 2. Extract speaker notes if present
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Speaker Notes]: {notes_text}")

        slide_full_text = "\n".join(slide_parts)
        if not slide_full_text:
            continue

        text_chunks = _chunk_text(slide_full_text, chunk_size, overlap)
        for chunk in text_chunks:
            chunks.append({
                "content": chunk,
                "page_number": slide_idx,
                "section": f"Slide {slide_idx}",
                "source": filename,
            })

    logger.info("Extracted %d chunks from PPTX '%s' across %d slides.", len(chunks), filename, len(prs.slides))
    return chunks
